"""PPTX 演示文稿解析器 — 输出 Markdown 格式（每页一个标题 + 内容）"""

import io
from typing import Any, Dict, List

from pptx import Presentation

from app.services.parser.base import BaseParser


class PPTXParser(BaseParser):
    """使用 python-pptx 提取 PPTX 内容，输出 Markdown 格式"""

    def parse(self, data: bytes, filename: str) -> Dict[str, Any]:
        prs = Presentation(io.BytesIO(data))
        parts: List[str] = []

        for i, slide in enumerate(prs.slides, 1):
            slide_parts: List[str] = []
            title_text = ""

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue
                        # 第一个非空文本框通常是标题
                        if not title_text and shape.shape_type == 13:  # PLACEHOLDER
                            title_text = text
                        else:
                            slide_parts.append(text)

                if shape.has_table:
                    table_md = _table_to_markdown(shape.table)
                    if table_md:
                        slide_parts.append(table_md)

            # 如果没有通过 shape_type 识别到标题，用第一段文本
            if not title_text and slide_parts:
                title_text = slide_parts.pop(0)

            # 组装这一页
            if title_text or slide_parts:
                page_md = f"## 第 {i} 页：{title_text}" if title_text else f"## 第 {i} 页"
                if slide_parts:
                    page_md += "\n\n" + "\n\n".join(slide_parts)
                parts.append(page_md)

        metadata: Dict[str, Any] = {"parse_method": "python-pptx", "output_format": "markdown"}

        return {
            "text": "\n\n".join(parts),
            "pages": len(prs.slides),
            "metadata": metadata,
        }


def _table_to_markdown(table) -> str:
    """将 PPTX 表格转为 Markdown 表格"""
    rows_data: List[List[str]] = []
    for row in table.rows:
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows_data.append(cells)

    if not rows_data:
        return ""

    header = rows_data[0]
    separator = ["---"] * len(header)
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(separator) + " |",
    ]
    for row in rows_data[1:]:
        while len(row) < len(header):
            row.append("")
        lines.append("| " + " | ".join(row[:len(header)]) + " |")

    return "\n".join(lines)
